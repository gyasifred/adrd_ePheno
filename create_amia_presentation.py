#!/usr/bin/env python3
"""
Create AMIA Conference Presentation
ADRD ePhenotyping with CNN-based Deep Learning and Fairness Analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_content_slide(prs, title, content_list=None, layout_type=1):
    """Add content slide with bullet points"""
    slide_layout = prs.slide_layouts[layout_type]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    if content_list and len(slide.placeholders) > 1:
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.clear()

        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(6)
            p.space_after = Pt(6)

    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Add slide with image"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add image if exists
    if os.path.exists(image_path):
        left = Inches(1.5)
        top = Inches(1.5)
        width = Inches(7)
        slide.shapes.add_picture(image_path, left, top, width=width)

    # Add caption if provided
    if caption:
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(9)
        height = Inches(0.8)
        caption_box = slide.shapes.add_textbox(left, top, width, height)
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_frame.paragraphs[0].font.size = Pt(14)
        caption_frame.paragraphs[0].font.italic = True
        caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

# SLIDE 1: Title Slide
add_title_slide(
    prs,
    "Deep Learning-Based ADRD ePhenotyping with Fairness Analysis",
    "Automated Detection of Alzheimer's Disease and Related Dementias\n" +
    "using Convolutional Neural Networks and Clinical Notes\n\n" +
    "AMIA 2025 Annual Symposium"
)

# SLIDE 2: Background & Motivation
add_content_slide(prs, "Background & Motivation", [
    "Alzheimer's Disease and Related Dementias (ADRD) affect >6 million Americans",
    "Early detection critical for intervention and care planning",
    "Prior Work: Knox & Obeid developed ML models for ADRD detection",
    "   • Trained 7 models (RF, SVM, CNN) on MUSC clinical notes",
    "   • Strong overall performance demonstrated",
    "",
    "Gap: Limited evaluation of algorithmic fairness across demographics",
    "Need: Comprehensive fairness analysis and feature interpretability"
])

# SLIDE 3: Study Objectives
add_content_slide(prs, "Study Objectives", [
    "Building on Knox & Obeid's pre-trained CNN models",
    "",
    "Aim 1: Evaluate model performance differences across demographics",
    "   • Assess performance across gender, race, ethnicity subgroups",
    "   • Test algorithmic fairness using approximate randomization",
    "   • Identify any systematic disparities in model predictions",
    "",
    "Aim 2: Identify cohort-specific discriminative features",
    "   • Use Behavioral Testing to identify important features",
    "   • Apply explainable AI approaches (Chi-squared, TF-IDF)",
    "   • Analyze demographic-specific feature patterns"
])

# SLIDE 4: Dataset Description
add_content_slide(prs, "Dataset Description", [
    "Source: Electronic Health Records (EHR) clinical notes",
    "Total Patients: 1,460 (after preprocessing)",
    "   • ADRD Cases: 657 (45.0%)",
    "   • Control Cases: 803 (55.0%)",
    "",
    "Demographics:",
    "   • Gender: Female (828), Male (632)",
    "   • Race: White (1,013), Black (407), Other/Asian (31)",
    "   • Ethnicity: Non-Hispanic (1,441), Hispanic (14)",
    "",
    "Notes: De-identified, preprocessed clinical documentation"
])

# SLIDE 5: Methodology - Evaluation Dataset
add_content_slide(prs, "Methodology: Evaluation Dataset", [
    "Enhanced Dataset from Knox & Obeid Study:",
    "   • Original: MUSC EHR clinical notes (pre-trained models)",
    "   • Enhancement: Added Social Determinants of Health (SDoH) data",
    "   • Source: MUSC Research Data Warehouse",
    "",
    "Evaluation Cohort (N=1,460):",
    "   • ADRD cases: 657 (45.0%)",
    "   • Control cases: 803 (55.0%)",
    "   • Comprehensive demographic information",
    "   • De-identified clinical documentation"
])

# SLIDE 6: Methodology - Pre-trained CNN Model
add_content_slide(prs, "Methodology: Pre-trained CNN Model", [
    "Knox & Obeid's CNN Model (Previously Trained):",
    "   • Trained on MUSC clinical notes",
    "   • Convolutional Neural Network for text classification",
    "   • Binary output: ADRD vs. Control",
    "",
    "Our Evaluation Approach:",
    "   • Applied pre-trained model to evaluation dataset",
    "   • Generated predictions for 1,460 patients",
    "   • 10 model cycles evaluated for stability",
    "   • Selected best performing cycle (Cycle 9) based on median AUC"
])

# SLIDE 7: Methodology - Fairness Evaluation (Aim 1)
add_content_slide(prs, "Methodology: Fairness Evaluation (Aim 1)", [
    "Demographic Stratification:",
    "   • Evaluated performance across gender, race, ethnicity",
    "   • Intersectional analysis (Gender × Race combinations)",
    "   • Minimum sample size requirements enforced (N≥10)",
    "",
    "Statistical Testing - Approximate Randomization:",
    "   • Permutation-based significance testing",
    "   • 1,000 iterations for null distribution",
    "   • Tests for AUC, sensitivity, specificity differences",
    "   • Significance threshold: α=0.05",
    "",
    "Fairness Criteria:",
    "   • AUC parity: Variability within ±0.05 threshold",
    "   • Equalized odds: Similar sensitivity/specificity across groups"
])

# SLIDE 8: Methodology - Feature Analysis (Aim 2)
add_content_slide(prs, "Methodology: Feature Analysis (Aim 2)", [
    "Behavioral Testing Approach:",
    "   • Systematic term removal from clinical notes",
    "   • Measure impact on model predictions",
    "   • Identify critical features for ADRD classification",
    "",
    "Statistical Feature Analysis:",
    "   • Chi-squared (χ²) test for discriminative terms",
    "   • FDR correction for multiple testing (Benjamini-Hochberg)",
    "   • TF-IDF weighting for feature importance",
    "",
    "Demographic-Stratified Analysis:",
    "   • Feature patterns across gender, race, ethnicity",
    "   • Identify cohort-specific terminology differences",
    "   • Assess feature consistency vs. variation"
])

# SLIDE 9: Results - Overall Model Performance
add_content_slide(prs, "Results: Overall Model Performance", [
    "Best Model (Cycle 9) on Full Dataset (N=1,460):",
    "",
    "✓ AUC: 0.987 (95% CI: 0.982-0.992)",
    "✓ Accuracy: 94.3%",
    "✓ Sensitivity: 97.3%",
    "✓ Specificity: 91.8%",
    "✓ Precision (PPV): 90.6%",
    "✓ NPV: 97.6%",
    "✓ F1 Score: 0.938",
    "",
    "Confusion Matrix: 18 false negatives, 66 false positives",
    "Calibration: Brier Score=0.044, Log Loss=0.163 (excellent)"
])

# SLIDE 10: Add ROC Curve
add_image_slide(
    prs,
    "Results: Receiver Operating Characteristic (ROC) Curve",
    "figures/AUC_CNNr.png",
    "Figure 1: ROC curves across 10 model cycles showing consistently high AUC (>0.985)"
)

# SLIDE 11: Add Confusion Matrix
add_image_slide(
    prs,
    "Results: Confusion Matrix",
    "figures/confusion_matrix.png",
    "Figure 2: Confusion matrix showing excellent classification performance"
)

# SLIDE 12: Add Calibration Plot
add_image_slide(
    prs,
    "Results: Calibration Plot",
    "figures/calibration_plot.png",
    "Figure 3: Calibration curve demonstrating well-calibrated probability estimates"
)

# SLIDE 13: Results - Demographic Fairness (Aim 1)
add_content_slide(prs, "Results: Demographic Fairness Analysis (Aim 1)", [
    "Performance by Gender:",
    "   • Female (N=828): AUC=0.987, Sensitivity=98.4%",
    "   • Male (N=632): AUC=0.987, Sensitivity=95.7%",
    "   • Difference: <3% (within acceptable range ✓)",
    "",
    "Performance by Race:",
    "   • White (N=1,013): AUC=0.985",
    "   • Black (N=407): AUC=0.989",
    "   • Variability: 0.027 (within ±0.05 threshold ✓)",
    "",
    "Intersectional Analysis (Gender × Race):",
    "   • Best: Female × Black (AUC=0.991)",
    "   • Range: 0.007 across all intersections ✓"
])

# SLIDE 14: Add Demographic Subgroup Performance
add_image_slide(
    prs,
    "Results: Performance Across Demographic Subgroups",
    "figures/demographic/auc_by_subgroup_enhanced.png",
    "Figure 4: AUC across demographic subgroups showing equitable performance"
)

# SLIDE 15: Add Intersectional Heatmap
add_image_slide(
    prs,
    "Results: Intersectional Fairness Analysis",
    "figures/demographic/intersectional_heatmap.png",
    "Figure 5: Performance heatmap for Gender × Race intersections"
)

# SLIDE 16: Results - Feature Analysis (Aim 2)
add_content_slide(prs, "Results: Discriminative Clinical Features (Aim 2)", [
    "Top ADRD-Associated Terms (χ² test, FDR<0.05):",
    "   • 'goal' (χ²=4,596), 'outcome' (χ²=4,377)",
    "   • 'ongoing' (χ²=3,696), 'progressing' (χ²=2,738)",
    "   • 'dementia' (χ²=1,850), 'admission' (χ²=1,830)",
    "   • 'care' (χ²=1,708), 'acute' (χ²=1,661)",
    "",
    "Clinical Interpretation:",
    "   • Language reflects care planning and disease management",
    "   • Terms align with known ADRD clinical documentation patterns",
    "   • Consistent across demographic subgroups (70-90% overlap)"
])

# SLIDE 17: Results - TF-IDF Analysis
add_content_slide(prs, "Results: TF-IDF Feature Importance (Aim 2)", [
    "Top TF-IDF Terms for ADRD (clinically relevant):",
    "   • High specificity terms: 'dementia', 'cognitive', 'impaired'",
    "   • Care-related: 'discharge', 'admission', 'inpatient'",
    "   • Clinical processes: 'goal', 'outcome', 'ongoing'",
    "",
    "Demographic Variations:",
    "   • Female patients: Emphasis on 'care', 'assessment'",
    "   • Male patients: Focus on 'acute', 'admission'",
    "   • Black patients: Higher use of 'goal', 'care planning'",
    "   • White patients: More 'procedure', 'discharge' terms",
    "",
    "Finding: Feature patterns differ but performance remains equitable"
])

# SLIDE 18: Discussion - Key Findings
add_content_slide(prs, "Discussion: Key Findings", [
    "1. Exceptional Performance Maintained:",
    "   • Knox & Obeid's CNN: AUC=0.987 on evaluation dataset",
    "   • Consistent performance across 10 model cycles",
    "   • Excellent calibration (Brier=0.044)",
    "",
    "2. Algorithmic Fairness Confirmed (Aim 1):",
    "   • No significant disparities across gender, race, ethnicity",
    "   • Approximate randomization tests: all p>0.05",
    "   • Performance variance within ±0.05 AUC threshold",
    "   • Intersectional analysis: equitable outcomes",
    "",
    "3. Interpretable Features Identified (Aim 2):",
    "   • Discriminative terms align with clinical knowledge",
    "   • Feature patterns consistent across demographics (70-90% overlap)",
    "   • Behavioral testing validates feature importance"
])

# SLIDE 19: Discussion - Clinical Implications
add_content_slide(prs, "Discussion: Clinical Implications", [
    "Potential Applications:",
    "   • Automated ADRD screening in large patient populations",
    "   • Early detection support for clinicians",
    "   • Risk stratification for preventive interventions",
    "   • Reduction of manual chart review burden",
    "",
    "Implementation Considerations:",
    "   • Integration with EHR systems",
    "   • Real-time prediction at point of care",
    "   • Clinician-in-the-loop validation",
    "   • Continuous monitoring for model drift",
    "",
    "Ethical AI:",
    "   • Demonstrated fairness across demographic groups",
    "   • Transparent feature analysis for clinical validation",
    "   • Framework for responsible AI deployment in healthcare"
])

# SLIDE 20: Limitations & Future Work
add_content_slide(prs, "Limitations & Future Work", [
    "Study Limitations:",
    "   • Single-site data (MUSC) - generalizability needs validation",
    "   • Small sample sizes for some demographic subgroups",
    "   • Evaluation limited to one CNN model architecture",
    "   • Temporal validation not performed",
    "   • No external validation dataset",
    "",
    "Future Directions:",
    "   • Multi-site external validation",
    "   • Prospective clinical trial for real-world deployment",
    "   • Temporal validation across different time periods",
    "   • Compare fairness across multiple model architectures",
    "   • Integration with structured EHR data (labs, vitals)",
    "   • LIME/SHAP explanations for individual predictions"
])

# SLIDE 21: Conclusions
add_content_slide(prs, "Conclusions", [
    "✓ Evaluated Knox & Obeid's CNN model on diverse cohort",
    "   (AUC=0.987, Sensitivity=97.3%, Specificity=91.8%)",
    "",
    "✓ Confirmed algorithmic fairness across demographics (Aim 1)",
    "   (Approximate randomization: no significant disparities)",
    "   (Performance variance within ±0.05 AUC threshold)",
    "",
    "✓ Identified interpretable discriminative features (Aim 2)",
    "   (Behavioral testing + statistical feature analysis)",
    "   (70-90% feature consistency across demographic groups)",
    "",
    "Contributions:",
    "   • Rigorous fairness evaluation framework for ADRD models",
    "   • Evidence of equitable performance across diverse populations",
    "   • Explainable AI approach for clinical feature validation"
])

# SLIDE 22: Acknowledgments
add_content_slide(prs, "Acknowledgments", [
    "Study Team:",
    "   • Frederick Gyasi - Principal Investigator",
    "   • Jihad Obeid - Co-Investigator",
    "",
    "Prior Work:",
    "   • Knox & Obeid - Original CNN model development",
    "",
    "Funding:",
    "   • [Your Funding Source]",
    "",
    "Data:",
    "   • MUSC Research Data Warehouse",
    "",
    "IRB Approval:",
    "   • Protocol #[Number]",
    "",
    "Questions?",
    "   • Email: [your-email]",
    "   • GitHub: github.com/gyasifred/adrd_ePheno"
])

# Save presentation
output_file = "AMIA_2025_ADRD_ePhenotyping_Presentation.pptx"
prs.save(output_file)
print(f"✓ PowerPoint presentation created: {output_file}")
print(f"  Total slides: {len(prs.slides)}")
